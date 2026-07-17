# core.py — roadmap_analyst (v4: visual Gantt by quarters)
#
# Reads .pptx roadmap files, extracts real work item names from text frames,
# builds a Russian reference roadmap template with:
#   1. Title slide
#   2. Program summary
#   3. WBS (all extracted works grouped by phase)
#   4. VISUAL Gantt chart (per-quarter timeline + chevron task bars per phase)
#   5. Checklist (PMI PMBOK 7 + ICB4)
#
# Path discipline: writes ONLY to a caller-supplied `out_path`. Plugin layer
# must pass `state_dir/reference_roadmap.pptx`. Direct-driver calls may also
# pass a /Users/.../Desktop/... path, but the Skill Review Checklist declares
# `fs` permission for that. The skill itself does not decide.
import json
import logging
import pathlib
import re
from collections import OrderedDict

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
    from pptx.enum.shapes import MSO_SHAPE
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

log = logging.getLogger("roadmap_analyst.core")


# =================== ШУМОВЫЕ ФИЛЬТРЫ ===================

MONTH_RX = re.compile(
    r'^\s*(?:Янв|Фев|Мар|Апр|Май|Июн|Июл|Авг|Сен|Окт|Ноя|Дек)[\s\u00A0]*\d{2}\s*$',
    re.I,
)
PERCENT_RX = re.compile(r'^\s*\d+\s*%\s*$')
DAY_MONTH_RX = re.compile(r'^\s*\d{1,2}[\.\s]\d{1,2}(?:[\.\s]\d{2,4})?\s*(?:\|.*)?$')
QUARTER_RX = re.compile(r'^\s*\d?q\d{2,4}\s*$', re.I)
TBD_RX = re.compile(r'^\s*TBD\s*$', re.I)
DATE_PREFIX_RX = re.compile(r'^\s*\d{1,2}\.\d{1,2}\s*\|')
LEGEND_RX = re.compile(
    r'(Просроч|Выполненн|Планов|под\s+риском|под\s+зависи|Перенос\s+вехи|'
    r'Веха\s+или\s+зависим|зависимость\s+под\s+риском)',
    re.I,
)

# Глагольные стартеры работ (расширенный словарь)
TASK_VERB_RX = re.compile(
    r'(?:^|\s)(Разработка|Разработке|Разработку|Разработки|'
    r'Подготовка|Подготовке|Подготовку|Подготовки|'
    r'Тестирование|Тестировани[юе]|Тестирования|'
    r'Ревизия|Ревзии?|Ревизию|'
    r'Создание|Создани[юе]|Создания|'
    r'Получение|Получени[юе]|Получения|'
    r'Настройка|Настройке|Настройку|Настройки|'
    r'Интеграция|Интеграци[юе]|Интеграции|'
    r'Миграция|Миграци[юе]|Миграции|'
    r'Анализ|Анализа|Анализу|Анализе|'
    r'Устранение|Устранени[юе]|Устранения|'
    r'Воссоздание|Воссоздани[юе]|Воссоздания|'
    r'Публикация|Публикаци[юе]|Публикации|'
    r'Синхронизация|Синхронизаци[юе]|Синхронизации|'
    r'Пилот|Пилота|Пилоту|'
    r'Подключение|Подключени[юе]|Подключения|'
    r'Отказ|Отказа|Отказу|'
    r'Расчёт|Расчет|Расчёта|Расчёту|'
    r'Доработка|Доработке|Доработку|Доработки|'
    r'Поставка|Поставке|Поставку|Поставки|'
    r'Вывод|Вывода|Выводу|'
    r'Выключение|Выключени[юе]|Выключения|'
    r'Разметка|Разметке|Разметку|Разметки|'
    r'Завершение|Завершени[юе]|Завершения|'
    r'Приёмка|Приёмк[аиу]|'
    r'Рассылка|Рассылк[аиу]|'
    r'Параллельн\w+\s+run|Parallel\s+run|'
    r'Аналитик\w+|Инсталляци[юя]|'
    r'Воссозда\w+|Отказ\s+от|'
    r'Сводн\w+\s+табли|'
    r'Экспорт\s+PDF|'
    r'Сопровождени[юе]|'
    r'Запуск\s+в\s+ПРОМ|вывод\s+в\s+ПРОМ)\b',
    re.I,
)

PROGRAM_TITLE_RX = re.compile(r'(Миграц\w*\s+с\s+Qlik\s+на\s+[^\\.]+)', re.I)


# =================== ИЗВЛЕЧЕНИЕ КОНТЕНТА ИЗ .PPTX ===================

def _clean_line(s: str) -> str:
    s = re.sub(r'[\u00A0\u2007\u202F]', ' ', s)
    s = re.sub(r'\s+', ' ', s).strip()
    return s


def parse_pptx_text(path: str):
    """Return list of dicts {slide, shape, text, raw} from a .pptx."""
    out = []
    prs = Presentation(path)
    for si, slide in enumerate(prs.slides, 1):
        for shi, sh in enumerate(slide.shapes):
            if not sh.has_text_frame:
                continue
            full = '\n'.join(
                ''.join(r.text for r in p.runs) for p in sh.text_frame.paragraphs
            ).strip()
            if not full:
                continue
            for line in full.split('\n'):
                line = _clean_line(line)
                if not line:
                    continue
                out.append({
                    'slide': si, 'shape': shi,
                    'text': line, 'raw_path': str(path),
                })
    return out


def _is_noise(line: str) -> bool:
    if not line or len(line) < 3:
        return True
    if MONTH_RX.match(line):
        return True
    if PERCENT_RX.match(line):
        return True
    if DAY_MONTH_RX.match(line) and not TASK_VERB_RX.search(line):
        return True
    if QUARTER_RX.match(line):
        return True
    if TBD_RX.match(line):
        return True
    if LEGEND_RX.search(line) and not TASK_VERB_RX.search(line):
        return True
    return False


def extract_work_items(parsed_rows):
    seen = OrderedDict()
    for row in parsed_rows:
        line = row['text']
        if _is_noise(line):
            continue
        if DATE_PREFIX_RX.match(line) and TASK_VERB_RX.search(line):
            line = DATE_PREFIX_RX.sub('', line, count=1).strip()
            if not line:
                continue
        if TASK_VERB_RX.search(line) or (len(line) >= 18 and len(line.split()) >= 2):
            key = line.lower()
            if key not in seen:
                seen[key] = {
                    'text': line,
                    'source_slide': row['slide'],
                    'source_path': pathlib.Path(row['raw_path']).name,
                }
    return list(seen.values())


def detect_program_name(parsed_rows):
    for row in parsed_rows:
        m = PROGRAM_TITLE_RX.search(row['text'])
        if m:
            name = _clean_line(m.group(1))
            name = re.sub(r'(\b\w+\b)\s+\1', r'\1', name)
            return name
    return 'Программа миграции BI'


def group_tasks_into_phases(items):
    """Buckets works into 5 PMI PMBOK 7 phases."""
    buckets = OrderedDict([
        ('Инициирование', []),
        ('Планирование', []),
        ('Исполнение', []),
        ('Контроль и тестирование', []),
        ('Завершение и ввод в эксплуатацию', []),
    ])
    dispatch = {
        'Инициирование': r'(Анализ|Пилот|Создание\s+ролево|Ревизи|Получение\s+доступ|'
                        r'Бизнес|кейс|initiat)',
        'Планирование': r'(Планир|Дизайн|архитектур|процесс|Разработка\s+архитектур|'
                       r'Синхронизац|разметк\w+\s+данных|Разработка\s+редактор|'
                       r'Анализ\s+источник)',
        'Исполнение': r'(Разработк|Подготовк|Интеграц|Устранени|'
                      r'Воссоздани|Подключени|Настройк|Поставк|Публикаци|'
                      r'Доработк|Создани|Инсталляци)',
        'Контроль и тестирование': r'(Тестирован|Parallel\s+run|Приёмк|'
                                  r'ПСИ|ПРОМ|вывод\s+в\s+ПРОМ|'
                                  r'Настройк\w+\s+на\s+ПСИ|Настройк\w+\s+на\s+ПРОМ|'
                                  r'Утвержде\w+|Расчёт\s+топ-?|Top\s*-\s*\d+)',
        'Завершение и ввод в эксплуатацию': r'(Запуск\s+в\s+ПРОМ|Вывод\s+в\s+ПРОМ|Отказ\s+от|'
                                            r'Выключен|завершен\s+пилот|'
                                            r'Супермаркет\s+данных|'
                                            r'Legacy\s+БД\s+выведен|'
                                            r'Дата-продукт|'
                                            r'Обратн\w+\s+связ|'
                                            r'Рассылк\s+pdf)',
    }
    phase_order = list(dispatch.keys())
    for item in items:
        text = item['text']
        placed = False
        for phase in phase_order:
            if re.search(dispatch[phase], text, re.I):
                buckets[phase].append(text)
                placed = True
                break
        if not placed:
            buckets['Исполнение'].append(text)
    return buckets


# =================== ВИЗУАЛЬНАЯ ДИАГРАММА ГАНТА ===================

# 8 кварталов = ~2 года; покрывает типичные BI-программы.
QUARTERS = [
    ('Q3 2026', RGBColor(0xC9, 0x4A, 0x55)),
    ('Q4 2026', RGBColor(0xB8, 0x40, 0x4A)),
    ('Q1 2027', RGBColor(0xA6, 0x36, 0x40)),
    ('Q2 2027', RGBColor(0x94, 0x2C, 0x36)),
    ('Q3 2027', RGBColor(0x82, 0x22, 0x2C)),
    ('Q4 2027', RGBColor(0x70, 0x18, 0x22)),
    ('Q1 2028', RGBColor(0x5E, 0x0E, 0x18)),
    ('Q2 2028', RGBColor(0x4C, 0x04, 0x0E)),
]

PHASE_COLORS = {
    'Инициирование': RGBColor(0x6A, 0x9F, 0xC9),
    'Планирование': RGBColor(0x66, 0xB2, 0x9B),
    'Исполнение': RGBColor(0xD5, 0x82, 0x4F),
    'Контроль и тестирование': RGBColor(0xC9, 0xA8, 0x66),
    'Завершение и ввод в эксплуатацию': RGBColor(0x91, 0x6C, 0xB2),
}

# Какой квартальный диапазон покрывает каждая фаза (для фонового chevron-а)
PHASE_QUARTER_SPAN = OrderedDict([
    ('Инициирование', (0, 1)),       # Q3-Q4 2026
    ('Планирование', (1, 2)),         # Q4 2026 - Q1 2027
    ('Исполнение', (2, 5)),           # Q1-Q4 2027
    ('Контроль и тестирование', (4, 6)),  # Q3 2027 - Q1 2028
    ('Завершение и ввод в эксплуатацию', (5, 7)),  # Q4 2027 - Q2 2028
])


def _hex_color(rgb):
    return '{:02X}{:02X}{:02X}'.format(rgb[0], rgb[1], rgb[2])


def _shape_text(shape, text, *, size=10, bold=False, color=None, align=None, italic=False):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.text = ''
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    else:
        p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    if color is not None:
        r.font.color.rgb = color


def _add_gantt_slide(prs, phases_data, title='Дорожная карта: работы по кварталам'):
    """Визуальная диаграмма Ганта: временная шкала + чевроны фаз + бары работ."""
    if not phases_data or not any(phases_data.values()):
        return False

    layout = prs.slide_layouts[6]
    s = prs.slides.add_slide(layout)

    # Slide title
    title_box = s.shapes.add_textbox(Inches(0.5), Inches(0.18), Inches(12.3), Inches(0.5))
    _shape_text(title_box, title, size=18, bold=True, color=RGBColor(0x40, 0x10, 0x18),
                align=PP_ALIGN.LEFT)

    # Геометрия
    PHASE_LABEL_W = 1.7      # Левая колонка для подписей фаз
    TIMELINE_LEFT = 2.25     # L начала временной шкалы
    TIMELINE_RIGHT = 13.05   # R конца временной шкалы
    TIMELINE_W = TIMELINE_RIGHT - TIMELINE_LEFT
    HEADER_TOP = 0.85        # Верх заголовков кварталов
    HEADER_H = 0.34          # Высота заголовка
    BODY_TOP = 1.30          # Верх тела диаграммы
    BODY_BOT = 6.65          # Низ тела (перед легендой)
    BODY_H = BODY_BOT - BODY_TOP
    LEGEND_TOP = 6.78        # Легенда снизу

    # Quarter width
    n_quarters = len(QUARTERS)
    quarter_w = TIMELINE_W / n_quarters

    # ===== Header row: quarter cells =====
    for qi, (qlabel, qcolor) in enumerate(QUARTERS):
        qbox = s.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(TIMELINE_LEFT + qi * quarter_w),
            Inches(HEADER_TOP),
            Inches(quarter_w),
            Inches(HEADER_H),
        )
        qbox.fill.solid()
        qbox.fill.fore_color.rgb = qcolor
        qbox.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        _shape_text(qbox, qlabel, size=11, bold=True,
                    color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)

    # Vertical quarter dividers (только в зоне тела)
    for qi in range(1, n_quarters):
        line = s.shapes.add_connector(
            1,  # STRAIGHT
            Inches(TIMELINE_LEFT + qi * quarter_w),
            Inches(HEADER_TOP + HEADER_H),
            Inches(TIMELINE_LEFT + qi * quarter_w),
            Inches(BODY_BOT),
        )
        line.line.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
        line.line.width = Pt(0.5)

    # ===== Phase rows =====
    phases_present = [(name, items) for name, items in phases_data.items() if items]
    n_phases = len(phases_present)
    if n_phases == 0:
        return False
    phase_band_h = BODY_H / n_phases

    for pi, (phase_name, items) in enumerate(phases_present):
        band_top = BODY_TOP + pi * phase_band_h
        band_centre = band_top + phase_band_h / 2
        pcolor = PHASE_COLORS.get(phase_name, RGBColor(0x66, 0x66, 0x66))

        # --- Phase label (rotated vertical text on left) ---
        plabel = s.shapes.add_textbox(
            Inches(0.4), Inches(band_top),
            Inches(PHASE_LABEL_W - 0.1), Inches(phase_band_h),
        )
        plabel.fill.solid()
        plabel.fill.fore_color.rgb = pcolor
        plabel.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        _shape_text(plabel, phase_name, size=10, bold=True,
                    color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.LEFT)

        # --- Phase background chevron (full-span on its quarter range) ---
        q_start, q_end = PHASE_QUARTER_SPAN.get(phase_name, (0, n_quarters - 1))
        chev_left_in = TIMELINE_LEFT + q_start * quarter_w
        chev_w_in = (q_end - q_start + 1) * quarter_w - 0.05
        bg = s.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(chev_left_in),
            Inches(band_top + 0.05),
            Inches(chev_w_in),
            Inches(phase_band_h - 0.10),
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = pcolor
        bg.line.fill.background()
        _shape_text(bg, f'  {phase_name}  ', size=10, bold=True,
                    color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.LEFT)

        # --- Task bars: top items spread within phase quarter range ---
        n_tasks = min(len(items), 8)   # макс. 8 полос на фазу
        if n_tasks == 0:
            continue
        task_band_h = (phase_band_h - 0.10) / max(n_tasks, 1)
        for ti in range(n_tasks):
            task_text = items[ti]
            # Distribute tasks across phase quarter range
            sub_q_start = q_start + (ti * (q_end - q_start)) // max(n_tasks, 1)
            sub_q_end = min(sub_q_start + max(1, (q_end - q_start + 1) // n_tasks), q_end)
            bar_left_in = TIMELINE_LEFT + sub_q_start * quarter_w + 0.06
            bar_w_in = (sub_q_end - sub_q_start + 1) * quarter_w - 0.12
            if bar_w_in < 0.6:
                bar_w_in = quarter_w - 0.12
            bar_top_in = band_top + 0.05 + ti * task_band_h + 0.03
            bar_h_in = max(0.12, task_band_h - 0.06)

            bar = s.shapes.add_shape(
                MSO_SHAPE.CHEVRON,
                Inches(bar_left_in), Inches(bar_top_in),
                Inches(bar_w_in), Inches(bar_h_in),
            )
            # Светлее фазы, чтобы было видно на фоне
            lighter = RGBColor(
                min(0xFF, int(pcolor[0]) + 0x30),
                min(0xFF, int(pcolor[1]) + 0x30),
                min(0xFF, int(pcolor[2]) + 0x30),
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = lighter
            bar.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            bar.line.width = Pt(0.5)
            # Truncate long task names
            disp = task_text if len(task_text) <= 48 else task_text[:45] + '…'
            _shape_text(bar, disp, size=8, color=RGBColor(0x20, 0x20, 0x20),
                        align=PP_ALIGN.LEFT)

            # Milestone diamond at end of bar (>)
            ms = s.shapes.add_shape(
                MSO_SHAPE.DIAMOND,
                Inches(bar_left_in + bar_w_in - 0.10), Inches(bar_top_in - 0.02),
                Inches(0.18), Inches(bar_h_in + 0.04),
            )
            ms.fill.solid()
            ms.fill.fore_color.rgb = pcolor
            ms.line.color.rgb = RGBColor(0x20, 0x20, 0x20)
            ms.line.width = Pt(0.4)

        # If there are more tasks than 8, show an overflow marker
        if len(items) > n_tasks:
            note = s.shapes.add_textbox(
                Inches(TIMELINE_LEFT + n_quarters * quarter_w - 1.4),
                Inches(band_top + 0.05),
                Inches(1.3), Inches(phase_band_h - 0.10),
            )
            _shape_text(note, f'+ ещё {len(items) - n_tasks} работ', size=8,
                        italic=True, color=RGBColor(0x60, 0x60, 0x60),
                        align=PP_ALIGN.LEFT)

    # ===== Bottom legend =====
    legend_box = s.shapes.add_textbox(
        Inches(0.5), Inches(LEGEND_TOP),
        Inches(12.3), Inches(0.4)
    )
    legend_tf = legend_box.text_frame
    legend_tf.word_wrap = True
    legend_tf.text = ''
    # Phase legend
    leg_p = legend_tf.paragraphs[0]
    leg_p.alignment = PP_ALIGN.LEFT
    leg_r = leg_p.add_run()
    leg_r.text = 'Фазы: '
    leg_r.font.size = Pt(10)
    leg_r.font.bold = True
    leg_r.font.color.rgb = RGBColor(0x20, 0x20, 0x20)
    for phase_name, _items in phases_present:
        r = leg_p.add_run()
        r.text = f'■ {phase_name}   '
        r.font.size = Pt(10)
        r.font.color.rgb = PHASE_COLORS.get(phase_name, RGBColor(0x80, 0x80, 0x80))

    # Milestone explanation paragraph
    note_p = legend_tf.add_paragraph()
    note_p.alignment = PP_ALIGN.LEFT
    note_r = note_p.add_run()
    note_r.text = '◆ — milestone (контрольная точка задачи). Шкала — календарные кварталы.'
    note_r.font.size = Pt(8)
    note_r.font.italic = True
    note_r.font.color.rgb = RGBColor(0x60, 0x60, 0x60)

    return True


# =================== LEGACY TEXT-BASED SLIDES ===================

def _set_text_frame(tf, lines, base_size=12):
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.text = ''
    for idx, ln in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = str(ln)
        if len(str(ln)) > 60:
            sz = max(8, base_size - 3)
        elif len(str(ln)) > 40:
            sz = max(9, base_size - 2)
        elif len(str(ln)) > 25:
            sz = max(10, base_size - 1)
        else:
            sz = base_size
        run.font.size = Pt(sz)


def _add_title_slide(prs, title, subtitle):
    layout = prs.slide_layouts[6]
    s = prs.slides.add_slide(layout)
    box = s.shapes.add_textbox(Inches(0.6), Inches(2.5), Inches(12.1), Inches(1.5))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.size = Pt(28)
    r.font.bold = True
    box2 = s.shapes.add_textbox(Inches(0.6), Inches(4.3), Inches(12.1), Inches(2.5))
    tf2 = box2.text_frame
    tf2.word_wrap = True
    tf2.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = subtitle
    r2.font.size = Pt(16)


def _add_content_slide(prs, title, bullets, base_size=14):
    layout = prs.slide_layouts[6]
    s = prs.slides.add_slide(layout)
    h = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.7))
    htf = h.text_frame
    htf.word_wrap = True
    htf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    hp = htf.paragraphs[0]
    hp.alignment = PP_ALIGN.LEFT
    hr = hp.add_run()
    hr.text = title
    hr.font.size = Pt(20)
    hr.font.bold = True
    body = s.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12.3), Inches(6.0))
    _set_text_frame(body.text_frame, bullets, base_size=base_size)


def build_reference_pptx(out_path: str, data: dict) -> bool:
    if not HAS_PPTX:
        return False
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    program_name = data.get('program_name', 'Программа миграции BI')
    project_count = data.get('project_count', 0)
    work_item_count = data.get('work_item_count', 0)
    phases = data.get('phases', OrderedDict())

    # 1. Title
    title_line = 'Эталонная дорожная карта программы'
    sub_lines = [
        program_name,
        f'Построено на {project_count} проектах / {work_item_count} уникальных работ',
        'Методическая основа: PMI PMBOK 7 + ICB4',
    ]
    _add_title_slide(prs, title_line, '\n'.join(sub_lines))

    # 2. Summary
    summary = [
        f'Программа: {program_name}',
        f'Количество проектов-источников: {project_count}',
        f'Извлечено уникальных работ: {work_item_count}',
        f'Распределение по фазам PMI PMBOK 7:',
    ]
    for phase, items in phases.items():
        summary.append(f'  - {phase}: {len(items)} работ')
    summary += [
        '',
        'Рекомендации по применению шаблона:',
        '  - Для нового проекта той же программы - взять фазы как структурный шаблон',
        '  - WBS - из работ внутри релевантной фазы; не выходить за её пределы без stage-gate',
        '  - Stage-gates выставлять между фазами',
    ]
    _add_content_slide(prs, 'Сводка по программе', summary, base_size=14)

    # 3. WBS (text-based)
    wbs_lines = []
    for phase, items in phases.items():
        if not items:
            continue
        wbs_lines.append(f'{phase}:')
        for it in items:
            wbs_lines.append(f'  - {it}')
        wbs_lines.append('')
    if not wbs_lines:
        wbs_lines = ['В исходных .pptx работы не обнаружены.']
    _add_content_slide(prs, 'WBS: работы по фазам PMI PMBOK 7', wbs_lines, base_size=12)

    # 4. VISUAL GANTT - новая фишка v4
    _add_gantt_slide(prs, dict(phases))

    # 5. Checklist
    checklist = [
        'Обязательные элементы дорожной карты (PMI PMBOK 7 / ICB4):',
        '  [✓] Инициирование: бизнес-кейс, заинтересованные стороны, scope',
        '  [✓] Планирование: архитектура, road-map фаз, ресурсы, бюджет, риски',
        '  [✓] Исполнение: декомпозированная WBS с конкретными работами',
        '  [✓] Контроль и тестирование: stage-gates между фазами, ПСИ/ПРОМ',
        '  [✓] Завершение: вывод в ПРОМ, документы передачи, отказ от legacy',
        '',
        'Методические рекомендации:',
        '  * Time-buffer 10-25% между крупными фазами',
        '  * Stage-gate после каждой фазы (контрольная точка с критерием приёмки)',
        '  * Зависимости между командами - явные',
        '  * Управление рисками: реестр рисков с владельцами',
    ]
    _add_content_slide(prs, 'Чек-лист: что должно быть в дорожной карте', checklist, base_size=13)

    out = pathlib.Path(out_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return True


# =================== ТОЧКА ВХОДА ===================

def handle_analyze_roadmaps(ctx, rm, prog='Программа миграции BI',
                            out_path=None):
    """Главная точка входа скилла.

    Параметры:
        ctx: dict-контекст со state_dir под ключом 'state_dir'.
        rm: list (или JSON-строка) путей к .pptx-файлам дорожных карт.
        prog: имя программы по умолчанию (если не найдено в исходниках).
        out_path: путь для .pptx-результата. Если None - кладём в state_dir.

    Возвращает: JSON-envelope со статусом, summary, списком файлов, метриками.

    Path discipline: пишет ТОЛЬКО в out_path. Desktop-копия намеренно
    отсутствует — это вопрос пользовательского workflow, не самой программы.
    """
    log.info('Analyzing roadmaps (v4: visual Gantt)')

    # Coerce input
    if isinstance(rm, str):
        try:
            rm = json.loads(rm)
        except Exception:
            rm = [rm]
    rm = rm or []
    if not isinstance(rm, (list, tuple)):
        rm = [rm]

    # Parse inputs
    parsed_all = []
    for item in rm:
        p = str(item)
        if not pathlib.Path(p).exists():
            log.warning('skip missing: %s', p)
            continue
        if p.lower().endswith('.pptx'):
            try:
                parsed_all.extend(parse_pptx_text(p))
            except Exception as e:
                log.warning('parse error %s: %s', p, e)
        elif p.lower().endswith('.pdf'):
            # Skip silently — текущая версия читает только .pptx
            continue

    if not parsed_all:
        return json.dumps({
            'status': 'error',
            'summary': 'Не удалось прочитать ни одного .pptx файла.',
            'out': [], 'msg': 'Нет входных данных.',
        }, ensure_ascii=False, indent=2)

    parsed_by_project = OrderedDict()
    for row in parsed_all:
        parsed_by_project.setdefault(row['raw_path'], []).append(row)
    project_count = len(parsed_by_project)

    work_items = extract_work_items(parsed_all)
    program_name = detect_program_name(parsed_all) or prog
    phases = group_tasks_into_phases(work_items)

    data = {
        'program_name': program_name,
        'project_count': project_count,
        'work_item_count': len(work_items),
        'phases': phases,
    }

    # Resolve output path (path confinement: caller-supplied)
    state_dir = (ctx or {}).get('state_dir') if isinstance(ctx, dict) else None
    if not out_path:
        if state_dir:
            out_path = str(pathlib.Path(state_dir) / 'reference_roadmap.pptx')
        else:
            return json.dumps({
                'status': 'error',
                'summary': 'Не задан out_path и state_dir недоступен.',
                'out': [], 'msg': 'Нет выходного пути.',
            }, ensure_ascii=False, indent=2)

    out_path = str(out_path)
    ok = build_reference_pptx(out_path, data)
    if not ok:
        return json.dumps({
            'status': 'error',
            'summary': 'PPTX-генерация недоступна (нет python-pptx).',
            'out': [], 'msg': 'Установите python-pptx.',
        }, ensure_ascii=False, indent=2)

    return json.dumps({
        'status': 'ok',
        'summary': (
            f'Анализ {project_count} проектов, извлечено {len(work_items)} работ. '
            f'Эталон сохранён: {out_path}'
        ),
        'program_name': program_name,
        'project_count': project_count,
        'work_items_total': len(work_items),
        'phases': {k: len(v) for k, v in phases.items()},
        'out': [out_path],
        'msg': 'Готово. Включает визуальный Gantt по кварталам.',
    }, ensure_ascii=False, indent=2)
